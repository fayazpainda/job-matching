from pptx import Presentation

prs = Presentation("Job_Matching_Presentation.pptx")

# Update Slide 6 (Architecture) — add UI explanation
prs.slides[5].notes_slide.notes_text_frame.text = """Let me break down the four modules that fulfill the assignment requirements.

Module 1, Data Input: We load 719 jobs from 71 real companies across 32 locations, and 100 candidate profiles with diverse skills — from Python backend to LLM engineering to robotics. All stored in structured CSV format.

Module 2, Analysis and Decision: Our scoring engine evaluates four dimensions. Skills match uses keyword overlap — what percentage of required skills does the candidate have? Experience match is proportional — if the job needs 5 years and the candidate has 3, that's 60%. Education match maps levels — Associate, Bachelor, Master, PhD. Location match gives 100% for exact match, 70% for remote, and 30% otherwise. All weights are normalized automatically.

Module 3, Results Display: This is our professional web dashboard. It shows stats cards at the top — top match score, jobs analyzed, average score, and number of strong matches. Below that, a ranked table with score bars and color-coded breakdown pills. Then interactive Chart.js bar charts. And an AI analysis panel with per-row Analyze buttons.

Module 4, Parameter Interaction: The web UI gives full control. A candidate dropdown, four weight sliders that auto-normalize, a location multi-select with Select All and Clear All buttons, minimum salary filter, and Top-N control up to 30 results.

ABOUT THE WEB INTERFACE:
The interface is built with Flask on the backend and pure HTML/CSS/JavaScript on the frontend — no frameworks needed. Let me describe what the user sees:

At the top is a clean navigation bar showing the app name "JobMatch AI", the total number of jobs and candidates, and whether the Claude API is connected or in demo mode.

Below that is the Controls Panel — a white card with three sections side by side:
- Left: A dropdown to select any of the 100 candidates. Once selected, a preview panel appears showing the candidate's self-introduction, skill tags in blue pills, expected salary, years of experience, and education level.
- Center: Four color-coded weight sliders — blue for Skills, green for Experience, purple for Education, amber for Location. Each slider goes from 0 to 100, and the system auto-normalizes them.
- Right: Location filters with Select All and Clear All buttons at the top, a multi-select list of all 32 cities, plus inputs for minimum salary and number of results.

Below the controls is a prominent blue "Find Best Matches" button with a search icon.

When clicked, the results appear in four sections:
1. Stats Row — four cards showing Top Match Score, Jobs Analyzed, Average Score, and Strong Matches count, each with an icon and large colored number.
2. Results Table — a professional table with columns for Rank (gold/silver/bronze badges), Position and Company, Location pill, Salary range, Score with a colored progress bar, Breakdown pills (Skills/Exp/Edu/Location scores), and an Analyze button on each row.
3. Charts — two side-by-side Chart.js horizontal bar charts. Left chart shows overall match scores. Right chart shows the per-dimension breakdown with a color legend.
4. AI Analysis Panel — a blue-tinted card with a lightning bolt icon. When you click Analyze on any row, this panel shows Claude's analysis with three sections: Match Highlights, Main Gaps, and Recommendation, rendered as formatted Markdown.

The entire interface uses a clean light theme with Inter font, subtle shadows, and smooth hover animations. It's responsive — works on both desktop and tablet screens."""

# Update Slide 8 (Demo) — add interface walkthrough
prs.slides[7].notes_slide.notes_text_frame.text = """Now let me show you the live demo. I'll walk through 6 steps and explain the interface as we go.

[OPEN BROWSER TO localhost:5000]

When the app loads, you see the navigation bar at the top showing "JobMatch AI" with "719 jobs" and "100 candidates" on the right. There's also a badge showing whether Claude API is connected or running in Demo Mode.

Step 1: I open the web app and select a candidate from the dropdown. As soon as I pick one — let's say a Python backend developer — a preview card appears below the dropdown showing their bio, skill tags as blue pills, expected salary, experience, and education. This gives you a quick profile overview before running the match.

Step 2: I click the blue "Find Best Matches" button. The page instantly populates with four sections:
- First, four stats cards appear at the top: the top match score (like 89.5), total jobs analyzed (719), the average score across all results, and how many are strong matches above 70.
- Below that, a ranked table shows the top results. Each row has a rank badge — gold for #1, silver for #2, bronze for #3. You can see the job title, company name in blue, location in a purple pill, salary range, a big score number with a colored progress bar, and breakdown pills showing individual dimension scores.

Step 3: Now I adjust the weight sliders. Watch the numbers update in real time — if I drag the Location slider higher and click Match again, jobs in the candidate's preferred city jump to the top.

Step 4: I use the Select All button to include every city, or I can hold Ctrl and pick specific cities from the list. The system filters through all 719 jobs instantly.

Step 5: I click the blue "Analyze" button on any table row. The AI Analysis panel at the bottom lights up with a loading animation, then displays Claude's structured response: Match Highlights as bullet points, Main Gaps, and a one-line Recommendation. This is rendered as Markdown with headers and formatting.

Step 6: I switch to a different candidate — an LLM Engineer. Click Match again, and the entire page re-renders with new rankings. Now LLM Engineer roles at OpenAI, Anthropic, and Cohere appear at the top.

This demonstrates the complete pipeline: interactive data input, quantitative 4-dimension scoring, professional visualization, and LLM-powered semantic analysis — all running locally in the browser."""

prs.save("Job_Matching_Presentation.pptx")
print("Interface explanation notes added!")
