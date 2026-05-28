from pptx import Presentation

prs = Presentation("Job_Matching_Presentation.pptx")

def fix_runs(slide, old_text, new_text):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

# Slide 3 — fix stat descriptions (run-level text may differ)
s3 = prs.slides[2]
fix_runs(s3, "test database", "71 real companies")
fix_runs(s3, "for matching", "diverse profiles")

# Slide 5 — fix Sort & Display box
s5 = prs.slides[4]
fix_runs(s5, "Ranked Table+ Dual Bar Charts+ Heatmap",
         "Stats Dashboard+ Ranked Table+ Dual Charts")
fix_runs(s5, "Heatmap", "Stats Dashboard")

# Slide 9 — fix duplicate limitation
s9 = prs.slides[8]
# Find and fix the duplicated "Skill matching" line
for shape in s9.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        full = "".join(r.text for r in para.runs)
        if "Skill matching is keyword-level" in full and "semantic" in full:
            if para.runs:
                para.runs[0].text = "No real-time job feed     → Connect to LinkedIn/Indeed APIs"
                for r in para.runs[1:]:
                    r.text = ""

prs.save("Job_Matching_Presentation.pptx")
print("Fixes applied!")
