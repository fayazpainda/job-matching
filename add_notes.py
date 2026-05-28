from pptx import Presentation

prs = Presentation("Job_Matching_Presentation.pptx")

notes = {
    0: """Good morning everyone. My name is Painda Mohammad Fayaz, student ID 22511204.

Today I'm presenting my final project for the Knowledge Engineering and Intelligent Decision-Making course — an LLM-Powered Intelligent Job Matching System.

This system helps HR teams and job seekers find the best candidate-job matches using a combination of rule-based scoring and Claude AI analysis. Let me walk you through the problem, the solution, and a live demo.""",

    1: """Here is the outline of my presentation. I'll cover 9 sections in about 10 minutes, then we'll have 5 minutes for questions.

We'll start with the problem background, look at the difficulties in current recruiting, then dive into the solution design and system architecture. I'll show how the LLM is integrated at every stage, demonstrate the key features with a live demo, evaluate the results, and finish with future work.

Let's begin with the problem.""",

    2: """So why do we need intelligent job matching?

Imagine you're an HR professional at a tech company. You receive hundreds of resumes for dozens of open positions. For each resume, you spend 5 to 10 minutes manually comparing skills, experience, education, and location preferences against job requirements. That's incredibly time-consuming.

Our system addresses this with a large-scale database: 719 job postings across 71 real tech companies — including Google, Meta, OpenAI, Anthropic, NVIDIA, Stripe, SpaceX — and 100 candidate profiles covering every major tech role. We score matches across 4 dimensions: Skills, Experience, Education, and Location.

The goal is to produce sortable, explainable matches with LLM-powered natural language explanations — so HR teams know not just WHO matches, but WHY.""",

    3: """Let me explain the four main pain points we're solving.

First, low efficiency. At scale, manual resume screening is simply not sustainable. HR teams can't afford to spend 5-10 minutes per resume when they have hundreds of applicants.

Second, subjective scoring. Different recruiters have different standards. One HR person might weigh skills heavily, another might prioritize location. There's no consistent, quantifiable standard.

Third, black-box rejections. Traditional ATS systems — Applicant Tracking Systems — give candidates a single number with no explanation. Nobody learns why they were rejected.

Fourth, multi-dimensional tradeoffs. Skills, experience, education, location, and salary all matter, but their relative importance changes depending on the role, the team, and the company. We need a flexible system that lets users adjust these weights.""",

    4: """Here is our solution design — a two-layer architecture combining rule-based scoring with LLM semantic understanding.

The data flows from left to right. First, we load job and candidate data from CSV files into our Flask web application. Then, a rule-based scoring engine evaluates each candidate-job pair across four dimensions using weighted scoring. The results are displayed in a professional web dashboard with ranked tables, dual bar charts, and stats cards. Finally, Claude AI generates natural language analysis for any selected match.

The key insight is at the bottom: the web interface provides live interaction — users can adjust weights, filter locations with Select All, and change the number of results in real time.

Why this combined approach? Pure rules are transparent but rigid — they can't understand that React and Vue are related skills. Pure LLM is flexible but expensive and unreliable. Our combined approach gives you the best of both: auditable quantitative scores plus semantic AI explanations.""",

    5: """Let me break down the four modules that fulfill the assignment requirements.

Module 1, Data Input: We load 719 jobs from 71 real companies across 32 locations, and 100 candidate profiles with diverse skills — from Python backend to LLM engineering to robotics. All stored in structured CSV format.

Module 2, Analysis and Decision: Our scoring engine evaluates four dimensions. Skills match uses keyword overlap — what percentage of required skills does the candidate have? Experience match is proportional — if the job needs 5 years and the candidate has 3, that's 60%. Education match maps levels — Associate, Bachelor, Master, PhD. Location match gives 100% for exact match, 70% for remote, and 30% otherwise. All weights are normalized automatically.

Module 3, Results Display: This is our professional web dashboard. It shows stats cards at the top — top match score, jobs analyzed, average score, and number of strong matches. Below that, a ranked table with score bars and color-coded breakdown pills. Then interactive Chart.js bar charts. And an AI analysis panel with per-row Analyze buttons.

Module 4, Parameter Interaction: The web UI gives full control. A candidate dropdown, four weight sliders that auto-normalize, a location multi-select with Select All and Clear All buttons, minimum salary filter, and Top-N control up to 30 results.""",

    6: """Now let me explain how the LLM — specifically Anthropic's Claude — enhances this system at every stage of development and at runtime.

During the design phase, I discussed with Claude how to decompose the problem into input, analysis, display, and recommendation modules. This helped me structure the project architecture.

For code generation, Claude helped me build the Flask web application UI and the Chart.js interactive visualizations. This saved significant development time.

During debugging, Claude helped me fix issues with weight normalization, CSV parsing edge cases, and font display problems.

At runtime — and this is the core feature — Claude acts as the decision-support engine. When a user clicks Analyze on any job match, Claude receives the candidate profile, job details, and system scores, and generates a structured analysis with three sections: Match Highlights, Main Gaps, and a Recommendation.

The prompt design is important: we set the system role as "an experienced and concise HR consultant," request Markdown output under 300 words, and limit tokens to 600 to control API costs. The model is configurable via environment variable — defaulting to Claude Sonnet for fast responses.

For optimization, Claude suggested giving location mismatches a score of 0.3 instead of 0, which preserves remote work options in the rankings.""",

    7: """Now let me show you the live demo. I'll walk through 6 steps.

[OPEN BROWSER TO localhost:5000]

Step 1: I open the web app and select a candidate from the dropdown — let's pick someone like a Python backend developer.

Step 2: I click Find Best Matches with the default weights. You can see the stats cards update — top match score, 719 jobs analyzed, and the results table fills with ranked matches showing score bars and breakdown pills.

Step 3: Now I'll adjust the weight sliders. Watch — if I increase the Location weight, jobs matching the candidate's preferred location jump to the top of the rankings.

Step 4: I can use Select All to include every city, or filter to specific locations. The system filters through all 719 jobs in real time.

Step 5: I click the Analyze button on any row — and Claude generates a detailed analysis with highlights, gaps, and a recommendation. This is the LLM in action.

Step 6: Now I switch to a different candidate — say, an LLM Engineer. The entire ranking re-sorts instantly, and the top matches change to LLM-related roles at companies like OpenAI, Anthropic, and Cohere.

This demonstrates the full pipeline: data input, quantitative scoring, interactive visualization, and LLM-powered explanation.""",

    8: """Let me summarize what the system can do and what could be improved.

On the capabilities side: We have quantitative 4-dimensional matching with adjustable weights. A professional web dashboard with stats cards, interactive charts, and score bars. Live interactivity with Select All, per-row AI analysis, and real-time filtering. Claude provides intelligent natural language explanations. And the system always works — even without an API key, it falls back to mock responses.

On the limitations side: First, skill matching is still keyword-level. The system doesn't know that React and Vue are related frontend frameworks. The fix would be to use embedding similarity for semantic matching. Second, we don't have a live job feed — connecting to LinkedIn or Indeed APIs would make the data always current. Third, each Claude API call takes 2-5 seconds — caching common match combinations would speed this up significantly. Fourth, there's no learning feedback loop yet — ideally, HR teams would rate matches, and the system would adjust weights using reinforcement learning.

The technology stack includes Python, Flask, HTML/CSS/JavaScript, Chart.js for visualizations, pandas and scikit-learn for data processing, the Anthropic Claude API for AI analysis, and marked.js for rendering Markdown responses.""",

    9: """Thank you for your attention.

To summarize: this project demonstrates a complete intelligent decision-support application that combines rule-based quantitative scoring with LLM semantic understanding. It covers all four required modules — data input, analysis and decision, results display, and parameter interaction — with 719 real job postings, 100 candidate profiles, and Claude AI-powered explanations.

The system runs locally as a web application — anyone can use it by opening a browser, no technical knowledge required.

I'm happy to answer any questions you have, or I can show additional features in the live demo.

Thank you!"""
}

for idx, text in notes.items():
    slide = prs.slides[idx]
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text.strip()

prs.save("Job_Matching_Presentation.pptx")
print("Speaker notes added to all 10 slides!")
