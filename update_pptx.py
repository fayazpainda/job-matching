from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from copy import deepcopy
import re

prs = Presentation("Job_Matching_Presentation.pptx")

def replace_in_slide(slide, replacements):
    """Replace text in all shapes of a slide, preserving formatting."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in para.runs)
            for old, new in replacements:
                if old in full_text:
                    full_text = full_text.replace(old, new)
            if para.runs:
                para.runs[0].text = full_text
                for run in para.runs[1:]:
                    run.text = ""

def replace_exact_run(slide, old_text, new_text):
    """Replace text in individual runs across all shapes."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

# =============================================================
# Slide 3 — Problem Background (update numbers)
# =============================================================
s3 = prs.slides[2]
replace_exact_run(s3, "80+", "719")
replace_exact_run(s3, "Jobs in ourtest database", "Jobs across71 companies")
replace_exact_run(s3, "25+", "100")
replace_exact_run(s3, "Candidate profilesfor matching", "Candidate profileswith diverse skills")

# =============================================================
# Slide 5 — Solution Design (update architecture)
# =============================================================
s5 = prs.slides[4]
replace_exact_run(s5, "CSV / UI", "CSV / Web UI")
replace_exact_run(s5, "Job & CandidateDatabases", "Job & CandidateDatabases (719 + 100)")
replace_exact_run(s5, "Ranked Table+ Dual Bar Charts+ Heatmap",
                  "Ranked Table+ Dual Bar Charts+ Stats Dashboard")
replace_exact_run(s5, "ipywidgets live interaction (weights / filters / Top-N)",
                  "Flask Web App — live interaction (weights / filters / Top-N / Select All)")

# =============================================================
# Slide 6 — System Architecture (update modules)
# =============================================================
s6 = prs.slides[5]
replace_exact_run(s6, "Load job database (77 jobs)", "Load job database (719 jobs)")
replace_exact_run(s6, "Load candidate database (26 candidates)", "Load candidate database (100 candidates)")
replace_exact_run(s6, "Interactive candid", "71 real companies, 32 locations, interactive candid")

replace_exact_run(s6, "Sorted results table with color gradient",
                  "Professional web dashboard with stats cards")
replace_exact_run(s6, "Dual bar charts (total + breakdown)",
                  "Sorted table with score bars + Dual bar charts")
replace_exact_run(s6, "Global candidate×job heatmap", "Per-row Analyze button for AI deep-dive")
replace_exact_run(s6, "Top-3 summary ", "Color-coded breakdown pills ")

replace_exact_run(s6, "9 ipywidgets controls:", "Web UI controls:")
replace_exact_run(s6, "4 weight sliders (auto-normalized)",
                  "4 weight sliders (auto-normalized)")
replace_exact_run(s6, "Location filter (multi-select)",
                  "Location filter (multi-select + Select All/Clear All)")
replace_exact_run(s6, "• Min ", "• Min salary + Top-N (up to 30) ")

# =============================================================
# Slide 7 — LLM Integration (minor updates)
# =============================================================
s7 = prs.slides[6]
replace_exact_run(s7, "LLM generated the ipywidgets UI template\nand the dual-view matplotlib chart",
                  "LLM generated the Flask web app UI\nand the Chart.js interactive visualizations")
replace_exact_run(s7, "LLM generated the ipywidgets UI templateand the dual-view matplotlib chart",
                  "LLM generated the Flask web app UI and the Chart.js interactive visualizations")

# =============================================================
# Slide 8 — Live Demo Flow (update for web app)
# =============================================================
s8 = prs.slides[7]
replace_exact_run(s8, "Select candidate Alex Johnson",
                  "Open web app at localhost:5000\nSelect candidate from dropdown")
replace_exact_run(s8, "(Python backend, 5 years experience)",
                  "(e.g. Kimberly Robinson, Python Backend)")
replace_exact_run(s8, "Run with default weights",
                  "Click Find Best Matches with default weights")
replace_exact_run(s8, "Top 5 shows Senior Python Backend",
                  "Top 10 shows matching jobs across")
replace_exact_run(s8, "Engineers at ~90 score",
                  "71 companies with score bars")
replace_exact_run(s8, "Raise the Location weight",
                  "Adjust weight sliders live")
replace_exact_run(s8, "Location-matched jobs jump to top",
                  "Watch rankings update instantly")
replace_exact_run(s8, "Filter to San Francisco only",
                  "Use Select All / filter specific cities")
replace_exact_run(s8, "Live filtering in action",
                  "Filter from 719 jobs in real time")
replace_exact_run(s8, "View Claude-generated analysis",
                  "Click Analyze on any job row")
replace_exact_run(s8, "Highlights / Gaps / Recommendation",
                  "AI generates Highlights / Gaps / Advice")
replace_exact_run(s8, "Switch to Harper Scott (LLM engineer)",
                  "Switch candidate (e.g. LLM Engineer profile)")
replace_exact_run(s8, "Top 1 auto-changes to LLM roles",
                  "Results re-rank to matching roles")

# =============================================================
# Slide 9 — Results & Evaluation (update capabilities + tech)
# =============================================================
s9 = prs.slides[8]
replace_exact_run(s9, "Clean visualization (dual bar charts + global heatmap)",
                  "Professional web dashboard (stats cards + charts + score bars)")
replace_exact_run(s9, "Live interactivity (9 ipywidgets controls)",
                  "Live interactivity (web UI with Select All, per-row AI analysis)")
replace_exact_run(s9, "Small sample size (~80 jobs, ~25 candidates)",
                  "Skill matching is keyword-level")
replace_exact_run(s9, "Connect to real job board APIs",
                  "Replace with embedding similarity (semantic)")
replace_exact_run(s9, "Skill matching is keyword-level     → Replace with embedding similarity",
                  "No real-time job feed     → Connect to LinkedIn/Indeed APIs")
replace_exact_run(s9,
    "Python  |  Jupyter Notebook  |  pandas  |  NumPy  |  matplotlib  |  scikit-learn  |  ipywidgets  |  Anthropic Claude API",
    "Python  |  Flask  |  HTML/CSS/JS  |  Chart.js  |  pandas  |  scikit-learn  |  Anthropic Claude API  |  marked.js")

prs.save("Job_Matching_Presentation.pptx")
print("PowerPoint updated successfully!")
