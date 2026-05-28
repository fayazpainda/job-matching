from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Light-theme color palette ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
BORDER_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MEDIUM_GRAY = RGBColor(0x6B, 0x7B, 0x8D)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
SUBTITLE_TEXT = RGBColor(0x4A, 0x5A, 0x6A)

PRIMARY = RGBColor(0x1A, 0x56, 0xDB)       # professional blue
PRIMARY_DARK = RGBColor(0x0F, 0x3D, 0xA8)  # darker blue for headers
PRIMARY_LIGHT = RGBColor(0xE8, 0xF0, 0xFE) # very light blue tint
TEAL = RGBColor(0x00, 0x89, 0x7B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GOLD = RGBColor(0xE6, 0x9F, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static", "dut_logo.png")

STUDENT_NAME = "PAINDA MOHAMMAD FAYAZ"
STUDENT_ID = "22511204"
COURSE_TITLE = "LLM-Assisted Intelligent Decision Application"
UNIVERSITY = "Dalian University of Technology"
DATE = "June 2026"
FOOTER_NOTE = f"Confidential – {STUDENT_NAME}"

TOTAL_SLIDES = 13


def set_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def rounded_rect(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def text(slide, left, top, width, height, content, size=18,
         color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font="Calibri",
         italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.font.italic = italic
    p.alignment = align
    return txBox


def bullets(text_frame, items, size=15, color=DARK_TEXT, spacing=Pt(6),
            bullet_color=None):
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0


def add_footer(slide, slide_num):
    line = rect(slide, Inches(0.6), Inches(6.85), Inches(12.133), Pt(0.75), BORDER_GRAY)

    text(slide, Inches(0.6), Inches(6.95), Inches(5), Inches(0.35),
         FOOTER_NOTE, size=8, color=MEDIUM_GRAY, italic=True)

    text(slide, Inches(5.5), Inches(6.95), Inches(4), Inches(0.35),
         COURSE_TITLE, size=8, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

    text(slide, Inches(10.5), Inches(6.95), Inches(2.233), Inches(0.35),
         f"{slide_num} / {TOTAL_SLIDES}", size=9, color=MEDIUM_GRAY,
         align=PP_ALIGN.RIGHT)


def add_logo(slide, left=Inches(12.0), top=Inches(0.2), height=Inches(0.5)):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, left, top, height=height)


def slide_header(slide, title_text, subtitle_text=None):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.15), PRIMARY_LIGHT)
    rect(slide, Inches(0), Inches(1.15), Inches(13.333), Pt(3), PRIMARY)

    text(slide, Inches(0.6), Inches(0.18), Inches(10.5), Inches(0.65),
         title_text, size=28, color=PRIMARY_DARK, bold=True)

    if subtitle_text:
        text(slide, Inches(0.6), Inches(0.7), Inches(10.5), Inches(0.35),
             subtitle_text, size=14, color=SUBTITLE_TEXT, italic=True)

    add_logo(slide)


def card(slide, left, top, width, height, title_text, items, accent_color=PRIMARY):
    rounded_rect(slide, left, top, width, height, WHITE, border_color=BORDER_GRAY)
    rect(slide, left, top, Inches(0.06), height, accent_color)

    text(slide, left + Inches(0.25), top + Inches(0.12), width - Inches(0.4), Inches(0.35),
         title_text, size=15, color=accent_color, bold=True)

    txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.5),
                                     width - Inches(0.4), height - Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    bullets(tf, items, size=12, color=SUBTITLE_TEXT, spacing=Pt(4))


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 1: TITLE SLIDE                                       ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)

rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), PRIMARY)
rect(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), PRIMARY)

rect(slide, Inches(0), Inches(2.0), Inches(13.333), Inches(3.6), PRIMARY_LIGHT)
rect(slide, Inches(0), Inches(2.0), Inches(13.333), Pt(2), PRIMARY)
rect(slide, Inches(0), Inches(5.6), Inches(13.333), Pt(2), PRIMARY)

if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.7), Inches(0.35), height=Inches(0.7))

text(slide, Inches(1), Inches(1.15), Inches(11.333), Inches(0.5),
     UNIVERSITY, size=16, color=PRIMARY, align=PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(2.4), Inches(11.333), Inches(0.9),
     "LLM-Powered Intelligent Job Matching System",
     size=36, color=PRIMARY_DARK, bold=True, align=PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(3.5), Inches(11.333), Inches(0.5),
     f"Final Project  —  {COURSE_TITLE}",
     size=18, color=SUBTITLE_TEXT, align=PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(0.04),
     "", size=1)
rect(slide, Inches(5.5), Inches(4.3), Inches(2.333), Pt(1.5), BORDER_GRAY)

text(slide, Inches(1), Inches(4.6), Inches(11.333), Inches(0.4),
     f"Presented by:  {STUDENT_NAME}", size=20, color=DARK_TEXT,
     bold=True, align=PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(5.1), Inches(11.333), Inches(0.35),
     f"Student ID: {STUDENT_ID}", size=15, color=SUBTITLE_TEXT, align=PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(5.9), Inches(11.333), Inches(0.35),
     DATE, size=14, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

text(slide, Inches(0.6), Inches(6.95), Inches(5), Inches(0.3),
     FOOTER_NOTE, size=8, color=MEDIUM_GRAY, italic=True)
text(slide, Inches(5.5), Inches(6.95), Inches(4), Inches(0.3),
     COURSE_TITLE, size=8, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Title Slide\n\n"
    "Good morning/afternoon, professors and classmates.\n"
    "My name is Painda Mohammad Fayaz, Student ID 22511204.\n"
    "Today I will present my final project for the LLM-Assisted Intelligent Decision "
    "Application course.\n"
    "My project is an LLM-Powered Intelligent Job Matching System that combines "
    "rule-based scoring with Claude AI to help HR teams match candidates to job postings.\n"
    "The presentation will take approximately 10 minutes, followed by 5 minutes of Q&A."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 2: OUTLINE                                           ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "Presentation Outline")
add_footer(slide, 2)

outline_items = [
    ("01", "Problem Background", "Why job matching matters"),
    ("02", "Current Difficulties", "Pain points in traditional recruiting"),
    ("03", "Solution Architecture", "Rule-based + LLM hybrid approach"),
    ("04", "Scoring Algorithm", "Four-dimension weighted matching"),
    ("05", "LLM Integration", "How Claude enhances the system"),
    ("06", "Interactive Interface", "Jupyter widgets & Flask web app"),
    ("07", "Live Demo Walkthrough", "End-to-end demonstration flow"),
    ("08", "Results & Capabilities", "What the system achieves"),
    ("09", "Limitations & Future Work", "Known gaps and roadmap"),
    ("10", "Summary & Q&A", "Key takeaways"),
]

for i, (num, title_str, desc) in enumerate(outline_items):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.5 + row * 1.05)

    num_shape = rounded_rect(slide, left, top + Inches(0.03), Inches(0.5), Inches(0.5), PRIMARY)
    text(slide, left + Inches(0.03), top + Inches(0.07), Inches(0.5), Inches(0.4),
         num, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(slide, left + Inches(0.65), top + Inches(0.0), Inches(4.8), Inches(0.3),
         title_str, size=17, color=DARK_TEXT, bold=True)
    text(slide, left + Inches(0.65), top + Inches(0.3), Inches(4.8), Inches(0.25),
         desc, size=12, color=MEDIUM_GRAY)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Outline\n\n"
    "Here is the roadmap for today’s presentation.\n"
    "I will start with the problem background, explain the difficulties in current "
    "recruiting processes, then walk through my solution architecture including the "
    "scoring algorithm and LLM integration.\n"
    "After that, I’ll show a live demo and discuss results, limitations, and future "
    "improvements.\nLet’s begin with the problem background."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 3: PROBLEM BACKGROUND                                ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "1. Problem Background", "Why intelligent job matching matters")
add_footer(slide, 3)

text(slide, Inches(0.8), Inches(1.45), Inches(7), Inches(0.4),
     "Scenario", size=20, color=PRIMARY_DARK, bold=True)
text(slide, Inches(0.8), Inches(1.9), Inches(11), Inches(0.8),
     "HR teams in tech companies need to match candidates to job postings for campus "
     "and social recruiting. The process is manual, time-consuming, and often "
     "inconsistent across different recruiters.",
     size=15, color=SUBTITLE_TEXT)

card(slide, Inches(0.8), Inches(3.0), Inches(3.7), Inches(1.5),
     "Target Users", [
         "• HR teams & recruiters",
         "• Job seekers & candidates",
         "• Recruiting platforms"
     ], PRIMARY)

card(slide, Inches(4.8), Inches(3.0), Inches(3.7), Inches(1.5),
     "Project Goal", [
         "• Sortable, explainable matches",
         "• Multi-dimensional scoring",
         "• AI-powered recommendations"
     ], TEAL)

card(slide, Inches(8.8), Inches(3.0), Inches(3.7), Inches(1.5),
     "Data Scale", [
         "• 77 job postings (LinkedIn-style)",
         "• 26 candidate profiles",
         "• 15+ tech categories covered"
     ], GOLD)

stat_bg = rounded_rect(slide, Inches(0.8), Inches(5.0), Inches(11.733), Inches(1.4),
                        PRIMARY_LIGHT, border_color=BORDER_GRAY)
stats = [
    ("5–10 min", "Per resume review\n(manual process)"),
    ("77 Jobs", "Across 10 US cities\n& remote positions"),
    ("26 Candidates", "Backend, Frontend,\nAI/ML, DevOps, etc."),
    ("4 Dimensions", "Skills, Experience,\nEducation, Location"),
]
for i, (val, desc) in enumerate(stats):
    x = Inches(1.2 + i * 3.0)
    text(slide, x, Inches(5.1), Inches(2.5), Inches(0.4),
         val, size=22, color=PRIMARY_DARK, bold=True, align=PP_ALIGN.CENTER)
    text(slide, x, Inches(5.5), Inches(2.5), Inches(0.5),
         desc, size=11, color=SUBTITLE_TEXT, align=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Problem Background\n\n"
    "The scenario is straightforward: HR teams in tech companies receive hundreds of "
    "resumes and need to match them against open job postings.\n"
    "Currently, each resume takes 5–10 minutes to manually compare against job "
    "requirements.\n"
    "Our dataset covers 77 jobs across major tech companies like Google, Meta, Amazon, "
    "OpenAI, and Anthropic, with 26 candidate profiles spanning backend, frontend, "
    "AI/ML, and DevOps roles.\n"
    "The system scores candidates across 4 dimensions: skills, experience, education, "
    "and location."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 4: CURRENT DIFFICULTIES                              ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "2. Current Difficulties",
             "Pain points in traditional recruiting workflows")
add_footer(slide, 4)

problems = [
    ("Low Efficiency",
     "HR spends 5–10 minutes per resume manually comparing qualifications "
     "to job requirements. This does not scale.",
     RED),
    ("Subjective Scoring",
     "Different recruiters apply different criteria. Results are inconsistent "
     "and hard to quantify or audit.",
     ORANGE),
    ("Black-Box Rejections",
     "Traditional ATS systems produce a single score with no explanation — "
     "candidates never learn why they were rejected.",
     PRIMARY_DARK),
    ("Multi-Dimensional Tradeoffs",
     "Skills, experience, education, location, and salary all matter, "
     "but the weights differ per role and per company.",
     TEAL),
]

for i, (title_str, desc, color) in enumerate(problems):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.5 + row * 2.65)

    rounded_rect(slide, left, top, Inches(5.8), Inches(2.3), WHITE,
                 border_color=BORDER_GRAY)
    rect(slide, left, top, Inches(0.07), Inches(2.3), color)

    num_badge = rounded_rect(slide, left + Inches(0.3), top + Inches(0.3),
                              Inches(0.5), Inches(0.5), color)
    text(slide, left + Inches(0.33), top + Inches(0.32), Inches(0.5), Inches(0.45),
         str(i + 1), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(slide, left + Inches(1.0), top + Inches(0.3), Inches(4.4), Inches(0.35),
         title_str, size=18, color=color, bold=True)
    text(slide, left + Inches(1.0), top + Inches(0.8), Inches(4.4), Inches(1.2),
         desc, size=14, color=SUBTITLE_TEXT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Current Difficulties\n\n"
    "There are four main pain points in traditional recruiting:\n"
    "1. Low efficiency – manual review takes 5–10 min per resume, doesn’t scale.\n"
    "2. Subjective scoring – different HR people, different criteria, inconsistent.\n"
    "3. Black-box rejections – ATS gives a score but no explanation.\n"
    "4. Multi-dimensional tradeoffs – skills, experience, education, location all "
    "matter differently per role.\n"
    "These problems motivate our hybrid rule + LLM approach."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 5: SOLUTION ARCHITECTURE                             ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "3. Solution Architecture",
             "Two-layer hybrid: Rule-Based Scoring + LLM Explanation")
add_footer(slide, 5)

pipeline_steps = [
    ("Data Input", "CSV / UI\nSelection", PRIMARY),
    ("Rule Scoring", "4-Dimension\nWeighted Score", TEAL),
    ("Sort & Display", "Table + Charts\n+ Heatmap", GREEN),
    ("LLM Explain", "Claude Analysis\n& Advice", GOLD),
]

for i, (title_str, desc, color) in enumerate(pipeline_steps):
    left = Inches(0.6 + i * 3.25)
    rounded_rect(slide, left, Inches(1.55), Inches(2.5), Inches(1.7), color)
    text(slide, left + Inches(0.1), Inches(1.65), Inches(2.3), Inches(0.4),
         title_str, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left + Inches(0.1), Inches(2.15), Inches(2.3), Inches(0.8),
         desc, size=14, color=WHITE, align=PP_ALIGN.CENTER)

    if i < 3:
        text(slide, left + Inches(2.55), Inches(2.0), Inches(0.6), Inches(0.5),
             "➡", size=24, color=PRIMARY, align=PP_ALIGN.CENTER)

text(slide, Inches(0.8), Inches(3.6), Inches(12), Inches(0.4),
     "Why This Design?", size=20, color=PRIMARY_DARK, bold=True)

comparisons = [
    ("Pure Rules", "Transparent & auditable\nbut rigid, no semantics", PRIMARY),
    ("Pure LLM", "Flexible & contextual\nbut opaque, expensive, unreliable", ORANGE),
    ("Our Hybrid  ✓", "Rules provide auditable scores;\nLLM adds semantic explanation",
     GREEN),
]

for i, (title_str, desc, color) in enumerate(comparisons):
    left = Inches(0.8 + i * 4.1)
    rounded_rect(slide, left, Inches(4.1), Inches(3.8), Inches(2.3), WHITE,
                 border_color=BORDER_GRAY)
    rect(slide, left, Inches(4.1), Inches(3.8), Pt(4), color)

    text(slide, left + Inches(0.3), Inches(4.3), Inches(3.2), Inches(0.35),
         title_str, size=17, color=color, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left + Inches(0.3), Inches(4.8), Inches(3.2), Inches(1.0),
         desc, size=14, color=SUBTITLE_TEXT, align=PP_ALIGN.CENTER)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Solution Architecture\n\n"
    "Our solution uses a two-layer combination:\n"
    "Layer 1: Rule-based scoring computes a weighted composite across 4 dimensions.\n"
    "Layer 2: LLM (Claude) generates natural-language explanations for the top match.\n\n"
    "Why not pure rules? Transparent but rigid — can’t understand semantic similarity.\n"
    "Why not pure LLM? Flexible but opaque, expensive, unreliable.\n"
    "Our hybrid gives auditable quantitative scores plus human-readable AI insights."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 6: SCORING ALGORITHM                                 ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "4. Scoring Algorithm",
             "Four-dimension weighted matching with adjustable parameters")
add_footer(slide, 6)

dimensions = [
    ("Skills Match", "50%",
     "|cand_skills ∩ job_skills|\n/ |job_skills|",
     "Keyword intersection ratio\nof candidate vs. required skills", PRIMARY),
    ("Experience", "20%",
     "1.0 if cand_exp ≥ job_min\nelse proportional",
     "Full score if meets minimum;\npartial credit otherwise", TEAL),
    ("Education", "15%",
     "Level mapping:\nAssociate=1, Bachelor=2,\nMaster=3, PhD=4",
     "1.0 if meets or exceeds\nrequired level", GREEN),
    ("Location", "15%",
     "1.0 = exact match\n0.7 = remote option\n0.3 = different city",
     "Preserves remote work\nand relocation options", GOLD),
]

for i, (name, weight, formula, desc, color) in enumerate(dimensions):
    left = Inches(0.6 + i * 3.15)
    top = Inches(1.45)

    rounded_rect(slide, left, top, Inches(2.95), Inches(5.1), WHITE,
                 border_color=BORDER_GRAY)

    rect(slide, left, top, Inches(2.95), Inches(0.85), color)
    text(slide, left + Inches(0.15), top + Inches(0.08), Inches(2.65), Inches(0.35),
         name, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left + Inches(0.15), top + Inches(0.42), Inches(2.65), Inches(0.3),
         f"Default weight: {weight}", size=12, color=RGBColor(0xE8, 0xF5, 0xE9),
         align=PP_ALIGN.CENTER)

    text(slide, left + Inches(0.2), top + Inches(1.05), Inches(2.55), Inches(0.25),
         "Formula", size=12, color=color, bold=True)

    rounded_rect(slide, left + Inches(0.12), top + Inches(1.35),
                 Inches(2.7), Inches(1.25), PRIMARY_LIGHT)
    text(slide, left + Inches(0.22), top + Inches(1.42), Inches(2.5), Inches(1.1),
         formula, size=12, color=DARK_TEXT)

    text(slide, left + Inches(0.2), top + Inches(2.8), Inches(2.55), Inches(0.25),
         "Explanation", size=12, color=color, bold=True)
    text(slide, left + Inches(0.2), top + Inches(3.1), Inches(2.55), Inches(1.2),
         desc, size=12, color=SUBTITLE_TEXT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Scoring Algorithm\n\n"
    "Each candidate–job pair is scored on four dimensions, normalized to 0–1:\n\n"
    "1. Skills Match (50%): intersection of candidate and required skills.\n"
    "2. Experience (20%): full score if meets minimum years; proportional otherwise.\n"
    "3. Education (15%): maps degrees to numeric levels.\n"
    "4. Location (15%): exact=1.0, remote=0.7, different city=0.3.\n\n"
    "The 0.3 for location mismatch was suggested by Claude to preserve remote options.\n"
    "All weights are user-adjustable and auto-normalized to sum to 1."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 7: LLM INTEGRATION                                   ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "5. LLM Integration",
             "Claude API for semantic analysis and natural-language recommendations")
add_footer(slide, 7)

stages_data = [
    ("Design Phase", "Discussed module decomposition\n(input / analysis / display / recommendation)"),
    ("Code Generation", "Generated ipywidgets UI template\nand matplotlib dual-view charts"),
    ("Debugging", "Helped debug font display, weight\nnormalization, CSV parsing issues"),
    ("Runtime (Core)", "Decision-support engine: generates\nhighlights / gaps / recommendation"),
    ("Optimization", "Suggested location mismatch score\nof 0.3 instead of 0 for remote options"),
]

for i, (stage, desc) in enumerate(stages_data):
    left = Inches(0.8)
    top = Inches(1.42 + i * 0.95)

    num_shape = rounded_rect(slide, left, top + Inches(0.08), Inches(0.4), Inches(0.4),
                              PRIMARY)
    text(slide, left + Inches(0.03), top + Inches(0.1), Inches(0.4), Inches(0.35),
         str(i + 1), size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(slide, left + Inches(0.55), top + Inches(0.0), Inches(2.5), Inches(0.3),
         stage, size=15, color=PRIMARY_DARK, bold=True)
    text(slide, left + Inches(0.55), top + Inches(0.3), Inches(4.5), Inches(0.55),
         desc, size=12, color=SUBTITLE_TEXT)

    if i < 4:
        rect(slide, left + Inches(0.18), top + Inches(0.52),
             Pt(2.5), Inches(0.45), PRIMARY_LIGHT)

prompt_bg = rounded_rect(slide, Inches(6.2), Inches(1.35), Inches(6.3), Inches(5.25),
                          WHITE, border_color=BORDER_GRAY)
rect(slide, Inches(6.2), Inches(1.35), Inches(6.3), Pt(4), PRIMARY)

text(slide, Inches(6.5), Inches(1.5), Inches(5.7), Inches(0.35),
     "Runtime Prompt Design", size=17, color=PRIMARY, bold=True)

prompt_items = [
    ("Role Setup", '"You are an experienced and concise HR consultant"'),
    ("Input Context", "Candidate profile + Job details + System scores"),
    ("Output Format", "3-section Markdown:\n   ✅ Match Highlights (2–3 bullets)\n"
                      "   ⚠️ Main Gaps (2–3 bullets)\n"
                      "   \U0001f4a1 Recommendation (1 sentence)"),
    ("Model", "claude-opus-4-7 (configurable via LLM_MODEL env var)"),
    ("Token Limit", "max_tokens: 600 (cost control)"),
    ("Cost Estimate", "~10 calls for demo, ~600 tokens each, ~$0.15 USD"),
]

for i, (label, value) in enumerate(prompt_items):
    y = Inches(2.05 + i * 0.73)
    text(slide, Inches(6.5), y, Inches(1.8), Inches(0.25),
         label, size=12, color=PRIMARY_DARK, bold=True)
    text(slide, Inches(8.3), y, Inches(3.9), Inches(0.65),
         value, size=11, color=SUBTITLE_TEXT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – LLM Integration\n\n"
    "Claude is used throughout the project lifecycle, not just at runtime:\n"
    "- Design: decomposed into input/analysis/display/recommendation modules\n"
    "- Code generation: ipywidgets UI and matplotlib charts\n"
    "- Debugging: font display, weight normalization, CSV parsing\n"
    "- Runtime core feature: natural-language match analysis\n"
    "- Optimization: suggested 0.3 location score\n\n"
    "The runtime prompt is carefully designed:\n"
    "- Role: ‘experienced HR consultant’\n"
    "- Structured 3-section output: highlights, gaps, recommendation\n"
    "- max_tokens=600 for cost control (~$0.15 total)"
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 8: INTERACTIVE INTERFACE                              ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "6. Interactive Interface",
             "Jupyter ipywidgets + Flask web application")
add_footer(slide, 8)

text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.35),
     "Jupyter Notebook (Primary)", size=17, color=PRIMARY_DARK, bold=True)

nb_card = rounded_rect(slide, Inches(0.8), Inches(1.85), Inches(5.5), Inches(4.5),
                        WHITE, border_color=BORDER_GRAY)
rect(slide, Inches(0.8), Inches(1.85), Inches(5.5), Pt(4), PRIMARY)

controls = [
    "•  Candidate selector dropdown (26 candidates)",
    "•  Skills weight slider (0.0 – 1.0)",
    "•  Experience weight slider (0.0 – 1.0)",
    "•  Education weight slider (0.0 – 1.0)",
    "•  Location weight slider (0.0 – 1.0)",
    "•  Location filter (multi-select, 10 cities + Remote)",
    "•  Minimum salary filter ($0 – $300K)",
    "•  Top-N control (3 – 15 results)",
    "•  LLM analysis toggle (enable/disable Claude)",
]

txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.0))
tf = txBox.text_frame
tf.word_wrap = True
bullets(tf, controls, size=13, color=DARK_TEXT, spacing=Pt(5))

text(slide, Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.35),
     "Flask Web App (Alternative)", size=17, color=PRIMARY_DARK, bold=True)

web_card = rounded_rect(slide, Inches(6.8), Inches(1.85), Inches(5.7), Inches(2.4),
                         WHITE, border_color=BORDER_GRAY)
rect(slide, Inches(6.8), Inches(1.85), Inches(5.7), Pt(4), TEAL)

web_features = [
    "•  Modern responsive web UI",
    "•  Real-time AJAX matching",
    "•  Interactive weight adjustment",
    "•  Location & salary filtering",
    "•  Score-colored result cards",
    "•  Claude AI analysis on demand",
    "•  Runs on localhost:5000",
]

txBox2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.2), Inches(2.2))
tf2 = txBox2.text_frame
tf2.word_wrap = True
bullets(tf2, web_features, size=13, color=DARK_TEXT, spacing=Pt(5))

text(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(0.35),
     "Technology Stack", size=16, color=PRIMARY_DARK, bold=True)

tech_card = rounded_rect(slide, Inches(6.8), Inches(4.95), Inches(5.7), Inches(1.8),
                          PRIMARY_LIGHT, border_color=BORDER_GRAY)

tech_items = [
    ("Python", "pandas, numpy, scikit-learn, matplotlib"),
    ("LLM", "Anthropic Claude API (claude-opus-4-7)"),
    ("UI", "ipywidgets + Flask + HTML/JS"),
    ("Data", "CSV (jobs.csv, candidates.csv)"),
]

for i, (tech, detail) in enumerate(tech_items):
    y = Inches(5.1 + i * 0.38)
    text(slide, Inches(7.1), y, Inches(1.2), Inches(0.3),
         tech, size=12, color=PRIMARY_DARK, bold=True)
    text(slide, Inches(8.3), y, Inches(3.9), Inches(0.3),
         detail, size=12, color=SUBTITLE_TEXT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Interactive Interface\n\n"
    "Two interfaces:\n"
    "1. Jupyter Notebook (primary): 9 interactive controls for candidate selection, "
    "weight adjustment, location filter, salary filter, and LLM toggle.\n"
    "2. Flask Web App (alternative): modern web UI on localhost:5000 with AJAX-powered "
    "real-time matching.\n\n"
    "Tech stack: Python (pandas, numpy, scikit-learn, matplotlib), Anthropic Claude API, "
    "ipywidgets + Flask."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 9: DEMO WALKTHROUGH                                  ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "7. Live Demo Walkthrough",
             "Recommended 3–4 minute demonstration flow")
add_footer(slide, 9)

demo_steps = [
    ("Step 1", "Select Candidate",
     'Pick "Alex Johnson"\n(Python backend, 5 years, Bachelor)',
     "Start with a well-rounded candidate", PRIMARY),
    ("Step 2", "Default Matching",
     "Run with default weights\n(Skills 50%, Exp 20%, Edu 15%, Loc 15%)",
     "Top 5 shows Senior Python Backend at ~90", TEAL),
    ("Step 3", "Adjust Weights", "Raise the Location weight slider",
     "Location-matched jobs jump to the top", GREEN),
    ("Step 4", "Apply Filters", 'Filter to "San Francisco, CA" only',
     "Live filtering narrows results instantly", GOLD),
    ("Step 5", "View AI Analysis", "Show Claude-generated explanation",
     "Highlights / Gaps / Recommendation output", PRIMARY),
    ("Step 6", "Switch Candidate", 'Change to "Harper Scott" (LLM engineer)',
     "Top 1 auto-changes to LLM Engineer roles", TEAL),
]

for i, (step, title_str, action, result, color) in enumerate(demo_steps):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.5 + row * 1.75)

    rounded_rect(slide, left, top, Inches(5.8), Inches(1.5), WHITE,
                 border_color=BORDER_GRAY)

    badge = rounded_rect(slide, left + Inches(0.15), top + Inches(0.12),
                          Inches(0.85), Inches(0.32), color)
    text(slide, left + Inches(0.18), top + Inches(0.12), Inches(0.8), Inches(0.3),
         step, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(slide, left + Inches(1.15), top + Inches(0.08), Inches(4.3), Inches(0.3),
         title_str, size=15, color=DARK_TEXT, bold=True)
    text(slide, left + Inches(1.15), top + Inches(0.42), Inches(4.3), Inches(0.5),
         action, size=11, color=SUBTITLE_TEXT)
    text(slide, left + Inches(1.15), top + Inches(1.05), Inches(4.3), Inches(0.3),
         f"→ {result}", size=11, color=color, bold=True)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Live Demo Walkthrough\n\n"
    "6-step demo flow (3–4 minutes):\n"
    "1. Pick Alex Johnson – Python backend, 5 years.\n"
    "2. Default weights – Senior Python Backend Engineer at ~90 score.\n"
    "3. Raise Location weight – location-matched jobs rise.\n"
    "4. Filter to San Francisco – live filtering.\n"
    "5. Show Claude AI analysis – highlights/gaps/recommendation.\n"
    "6. Switch to Harper Scott – top results change to LLM roles.\n\n"
    "This showcases scoring, weight adjustment, filtering, visualization, and LLM analysis."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 10: RESULTS & CAPABILITIES                           ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "8. Results & Capabilities", "What the system achieves")
add_footer(slide, 10)

capabilities = [
    ("Quantitative\nMatching", "4-dimension weighted scoring\nwith auto-normalization",
     PRIMARY),
    ("Clean\nVisualization", "Dual bar charts (total +\nbreakdown) + global heatmap",
     TEAL),
    ("Live\nInteractivity", "9 controls with instant\nupdates and filtering",
     GREEN),
    ("AI\nExplanation", "Natural-language recommendation\nrationale via Claude",
     GOLD),
]

for i, (title_str, desc, color) in enumerate(capabilities):
    left = Inches(0.6 + i * 3.15)
    rounded_rect(slide, left, Inches(1.45), Inches(2.95), Inches(2.15), WHITE,
                 border_color=BORDER_GRAY)
    rect(slide, left, Inches(1.45), Inches(2.95), Inches(0.7), color)
    text(slide, left + Inches(0.15), Inches(1.5), Inches(2.65), Inches(0.6),
         title_str, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left + Inches(0.15), Inches(2.3), Inches(2.65), Inches(1.0),
         desc, size=13, color=SUBTITLE_TEXT, align=PP_ALIGN.CENTER)

text(slide, Inches(0.8), Inches(3.95), Inches(11.5), Inches(0.35),
     "Assignment Requirements Checklist", size=18, color=PRIMARY_DARK, bold=True)

rounded_rect(slide, Inches(0.8), Inches(4.4), Inches(11.733), Inches(2.3),
             WHITE, border_color=BORDER_GRAY)

requirements = [
    ("Complete input → analysis → display → recommendation app",
     "✅ All 4 modules implemented"),
    ("Python + Jupyter Notebook", "✅ job_matching.ipynb"),
    ("Use pandas / matplotlib / scikit-learn / ipywidgets",
     "✅ All present and functional"),
    ("Data input, analysis, display, parameter interaction",
     "✅ Modules 1–4 with 9 controls"),
    ("Not chat logs, static docs, or non-runnable code",
     "✅ Fully runnable end-to-end"),
]

for i, (req, status) in enumerate(requirements):
    y = Inches(4.5 + i * 0.4)
    if i % 2 == 0:
        rounded_rect(slide, Inches(0.85), y, Inches(11.633), Inches(0.38),
                      PRIMARY_LIGHT)
    text(slide, Inches(1.0), y + Inches(0.06), Inches(7.0), Inches(0.28),
         req, size=12, color=DARK_TEXT)
    text(slide, Inches(8.2), y + Inches(0.06), Inches(4.0), Inches(0.28),
         status, size=12, color=GREEN, bold=True)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Results & Capabilities\n\n"
    "Four key capabilities:\n"
    "1. Quantitative matching across 4 dimensions with adjustable weights\n"
    "2. Dual-view visualization: bar charts and global heatmap\n"
    "3. Live interactivity with 9 controls\n"
    "4. AI-powered explanations via Claude\n\n"
    "All five assignment requirements are met — complete pipeline, Python + Jupyter, "
    "required libraries, all four modules, and fully runnable code."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 11: LIMITATIONS & FUTURE WORK                        ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "9. Limitations & Future Work",
             "Known gaps and improvement roadmap")
add_footer(slide, 11)

text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.35),
     "Current Limitations", size=17, color=RED, bold=True)

limitations = [
    ("Keyword-Level Matching",
     "Skills matching doesn’t recognize\nReact ↔ Vue as related frameworks"),
    ("Limited Sample Size",
     "Only ~80 jobs and ~25 candidates;\nnot validated against real resumes"),
    ("LLM Latency & Cost",
     "Each API call takes 2–5 seconds;\nrequires internet connection"),
    ("No Learning Feedback",
     "System doesn’t learn from HR\nratings of recommendations"),
]

for i, (title_str, desc) in enumerate(limitations):
    top = Inches(1.85 + i * 1.15)
    rounded_rect(slide, Inches(0.8), top, Inches(5.5), Inches(0.95), WHITE,
                 border_color=BORDER_GRAY)
    rect(slide, Inches(0.8), top, Inches(0.06), Inches(0.95), RED)
    text(slide, Inches(1.1), top + Inches(0.08), Inches(4.8), Inches(0.25),
         f"⚠  {title_str}", size=13, color=RED, bold=True)
    text(slide, Inches(1.1), top + Inches(0.35), Inches(4.8), Inches(0.55),
         desc, size=11, color=SUBTITLE_TEXT)

text(slide, Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.35),
     "Future Improvements", size=17, color=GREEN, bold=True)

improvements = [
    ("Semantic Embeddings",
     "Replace keyword sets with\nembedding-model similarity"),
    ("Resume Parsing",
     "Upload PDF/Word resumes;\nLLM auto-extracts structure"),
    ("Bidirectional Matching",
     "Also recommend candidates\nfor a given job posting"),
    ("Reinforcement Learning",
     "HR rates recommendations;\nweights auto-adjust via RL"),
    ("Multi-Model Voting",
     "Call multiple LLMs, aggregate\nfor more reliable results"),
]

for i, (title_str, desc) in enumerate(improvements):
    top = Inches(1.85 + i * 0.95)
    rounded_rect(slide, Inches(6.8), top, Inches(5.7), Inches(0.78), WHITE,
                 border_color=BORDER_GRAY)
    rect(slide, Inches(6.8), top, Inches(0.06), Inches(0.78), GREEN)

    num_badge = rounded_rect(slide, Inches(7.0), top + Inches(0.18),
                              Inches(0.32), Inches(0.32), GREEN)
    text(slide, Inches(7.01), top + Inches(0.18), Inches(0.32), Inches(0.3),
         str(i + 1), size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, Inches(7.5), top + Inches(0.06), Inches(2.0), Inches(0.25),
         title_str, size=12, color=GREEN, bold=True)
    text(slide, Inches(9.5), top + Inches(0.06), Inches(2.8), Inches(0.65),
         desc, size=11, color=SUBTITLE_TEXT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Limitations & Future Work\n\n"
    "Limitations:\n"
    "1. Keyword matching – doesn’t understand React ≈ Vue. Solution: embeddings.\n"
    "2. Small sample – 80 jobs, 25 candidates. Needs real databases.\n"
    "3. LLM latency – 2–5 sec per call. Could cache common combinations.\n"
    "4. No learning loop – doesn’t improve from HR feedback.\n\n"
    "Future work: semantic embeddings, resume parsing, bidirectional matching, "
    "RL from feedback, multi-model voting."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 12: SUMMARY                                          ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, OFF_WHITE)
slide_header(slide, "10. Summary", "Key takeaways from the project")
add_footer(slide, 12)

takeaways = [
    ("Hybrid Architecture",
     "Rule-based scoring provides transparency and auditability;\n"
     "LLM adds semantic understanding and natural-language explanations.\n"
     "Neither alone is sufficient — the combination delivers the best results.",
     PRIMARY),
    ("Full-Stack Implementation",
     "Complete pipeline from data input through analysis, visualization,\n"
     "and AI-powered recommendations. Two interfaces: Jupyter + Flask web app.\n"
     "All assignment requirements met with runnable, end-to-end code.",
     TEAL),
    ("Practical LLM Usage",
     "Claude is used across the full project lifecycle: design, coding,\n"
     "debugging, runtime analysis, and optimization. The runtime prompt\n"
     "is carefully engineered for structured, cost-effective output.",
     GREEN),
]

for i, (title_str, desc, color) in enumerate(takeaways):
    top = Inches(1.45 + i * 1.75)
    rounded_rect(slide, Inches(0.8), top, Inches(11.733), Inches(1.5), WHITE,
                 border_color=BORDER_GRAY)
    rect(slide, Inches(0.8), top, Inches(0.07), Inches(1.5), color)

    num_badge = rounded_rect(slide, Inches(1.1), top + Inches(0.3),
                              Inches(0.55), Inches(0.55), color)
    text(slide, Inches(1.13), top + Inches(0.32), Inches(0.52), Inches(0.5),
         str(i + 1), size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text(slide, Inches(1.9), top + Inches(0.12), Inches(10.0), Inches(0.35),
         title_str, size=19, color=color, bold=True)
    text(slide, Inches(1.9), top + Inches(0.5), Inches(10.0), Inches(0.9),
         desc, size=14, color=SUBTITLE_TEXT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Summary\n\n"
    "Three key takeaways:\n\n"
    "1. Hybrid rule + LLM architecture is the right approach for decision-support. "
    "Rules give auditability; LLM gives semantic understanding.\n\n"
    "2. Full-stack implementation covering data input, analysis, visualization, "
    "and AI recommendations with Jupyter and web interfaces.\n\n"
    "3. LLM used throughout the entire lifecycle — design through optimization. "
    "Runtime prompt engineered for structured output.\n\n"
    "Thank you for your attention. I’m happy to answer any questions."
)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 13: THANK YOU / Q&A                                  ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)

rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), PRIMARY)
rect(slide, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12), PRIMARY)

rect(slide, Inches(0), Inches(2.0), Inches(13.333), Inches(3.8), PRIMARY_LIGHT)
rect(slide, Inches(0), Inches(2.0), Inches(13.333), Pt(2), PRIMARY)
rect(slide, Inches(0), Inches(5.8), Inches(13.333), Pt(2), PRIMARY)

if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.7), Inches(0.4), height=Inches(0.65))

text(slide, Inches(1), Inches(1.15), Inches(11.333), Inches(0.45),
     UNIVERSITY, size=15, color=PRIMARY, align=PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(2.35), Inches(11.333), Inches(0.7),
     "Thank You!", size=44, color=PRIMARY_DARK, bold=True,
     align=PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(3.2), Inches(11.333), Inches(0.45),
     "Questions & Discussion", size=20, color=SUBTITLE_TEXT,
     align=PP_ALIGN.CENTER)

rect(slide, Inches(5.5), Inches(3.9), Inches(2.333), Pt(1.5), BORDER_GRAY)

info_items = [
    f"Student:  {STUDENT_NAME}",
    f"Student ID:  {STUDENT_ID}",
    f"Course:  {COURSE_TITLE}",
    f"Email:  fayazpainda@mail.dlut.edu.cn",
]

for i, item in enumerate(info_items):
    text(slide, Inches(4), Inches(4.2 + i * 0.42), Inches(5.333), Inches(0.35),
         item, size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)

text(slide, Inches(0.6), Inches(6.95), Inches(5), Inches(0.3),
     FOOTER_NOTE, size=8, color=MEDIUM_GRAY, italic=True)
text(slide, Inches(5.5), Inches(6.95), Inches(4), Inches(0.3),
     COURSE_TITLE, size=8, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
text(slide, Inches(10.5), Inches(6.95), Inches(2.233), Inches(0.3),
     f"{TOTAL_SLIDES} / {TOTAL_SLIDES}", size=9, color=MEDIUM_GRAY,
     align=PP_ALIGN.RIGHT)

notes = slide.notes_slide
notes.notes_text_frame.text = (
    "PRESENTER NOTES – Thank You & Q&A\n\n"
    "Thank you for your attention!\n"
    "I’m happy to answer any questions about the system architecture, "
    "the scoring algorithm, the LLM integration, or the demo.\n\n"
    "To try the system:\n"
    "1. Clone the project folder\n"
    "2. pip install -r requirements.txt\n"
    "3. Add your Anthropic API key to .env\n"
    "4. jupyter notebook job_matching.ipynb\n"
    "5. Or: python app.py (web interface at localhost:5000)\n\n"
    "Contact: fayazpainda@mail.dlut.edu.cn"
)


# ── Save ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "Final_Project_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {TOTAL_SLIDES}")
