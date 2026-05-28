from pptx import Presentation

prs = Presentation("Job_Matching_Presentation.pptx")

def fix_para_text(slide, old_fragment, new_full_text):
    """Find a paragraph containing old_fragment, rewrite its first run, clear the rest."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if old_fragment in full:
                if para.runs:
                    para.runs[0].text = new_full_text
                    for r in para.runs[1:]:
                        r.text = ""
                return True
    return False

# Slide 3 — fix stat card descriptions
s3 = prs.slides[2]
fix_para_text(s3, "Jobs in our71 real companies", "Jobs across\n71 real companies")
fix_para_text(s3, "Candidate profilesdiverse profiles", "Candidate profiles\nwith diverse skills")

# Slide 9 — fix the limitation that lost its emoji
s9 = prs.slides[8]
fix_para_text(s9, "No real-time job feed", "⚠️  No real-time job feed     → Connect to LinkedIn/Indeed APIs")

prs.save("Job_Matching_Presentation.pptx")
print("Final fixes applied!")
